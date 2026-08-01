"""Generate the Sonata Knowledge Assistant (SKA) Phase-1 budget-proposal deck.

Run with the POC venv python (python-pptx is installed there):
    "C:\\Users\\atgupta\\sonata-kb\\03-poc\\agent\\venv\\Scripts\\python.exe" build_deck.py
Output: 00-overview/Sonata-Knowledge-Assistant-Phase1-Proposal.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
INK      = RGBColor(0x0B, 0x0B, 0x0B)   # primary ink
INK2     = RGBColor(0x52, 0x51, 0x4E)   # secondary ink
MUTED    = RGBColor(0x89, 0x87, 0x81)   # muted
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PAGE     = RGBColor(0xFC, 0xFC, 0xFB)   # slide background
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
HAIR     = RGBColor(0xE1, 0xE0, 0xD9)   # hairline border
NAVY     = RGBColor(0x10, 0x42, 0x81)   # brand deep blue
BLUE     = RGBColor(0x2A, 0x78, 0xD6)   # accent blue
BLUE_BG  = RGBColor(0xE4, 0xEE, 0xFB)   # light blue fill
ORANGE   = RGBColor(0xEB, 0x68, 0x34)
AQUA     = RGBColor(0x1B, 0xAF, 0x7A)
GOLD     = RGBColor(0xED, 0xA1, 0x00)
GOOD     = RGBColor(0x0C, 0xA3, 0x0C)
GRAY     = RGBColor(0xF2, 0xF3, 0xF5)

FONT = "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers
def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=PAGE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1, round_=False, radius=0.08):
    typ = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.line_spacing = para.get("line_spacing", 1.06)
        p.space_after = Pt(para.get("space_after", 0))
        for run in para["runs"]:
            r = p.add_run(); r.text = run["t"]
            f = r.font
            f.size = Pt(run.get("size", 16)); f.bold = run.get("bold", False)
            f.color.rgb = run.get("color", INK); f.name = run.get("font", FONT)
    return tb

def bullets(slide, x, y, w, h, items, size=15, gap=8, marker=BLUE):
    paras = []
    for item in items:
        if isinstance(item, str):
            runs = [{"t": item, "size": size, "color": INK}]
        elif isinstance(item, tuple):
            txt, ov = item
            runs = [{"t": txt, "size": ov.get("size", size), "bold": ov.get("bold", False),
                     "color": ov.get("color", INK)}]
        else:  # list of run dicts
            runs = item
        paras.append({"runs": [{"t": "▪  ", "size": size, "bold": True, "color": marker}] + runs,
                      "space_after": gap, "line_spacing": 1.05})
    return text(slide, x, y, w, h, paras)

def header(slide, title, num, kicker=None):
    bg(slide)
    rect(slide, 0, 0, SW, 0.16, NAVY)
    add_title = [{"runs": [{"t": title, "size": 25, "bold": True, "color": NAVY}]}]
    if kicker:
        add_title.append({"runs": [{"t": kicker, "size": 12.5, "color": INK2}], "space_after": 0})
    text(slide, 0.5, 0.42, 11.5, 1.0, add_title)
    rect(slide, 0.53, 1.28, 2.6, 0.035, BLUE)
    # footer
    text(slide, 0.5, SH - 0.42, 8.0, 0.3,
         [{"runs": [{"t": "Sonata Knowledge Assistant · Phase 1 budget proposal · Aug 2026",
                    "size": 9, "color": MUTED}]}])
    text(slide, SW - 1.0, SH - 0.42, 0.6, 0.3,
         [{"runs": [{"t": str(num), "size": 10, "color": MUTED}], "align": PP_ALIGN.RIGHT}])

def stat_tile(slide, x, y, w, h, value, label, accent=BLUE, value_color=None, value_size=32):
    rect(slide, x, y, w, h, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.07)
    rect(slide, x, y, w, 0.11, accent)
    text(slide, x + 0.22, y + 0.32, w - 0.44, h - 0.5, [
        {"runs": [{"t": value, "size": value_size, "bold": True, "color": value_color or accent}]},
        {"runs": [{"t": label, "size": 12.5, "color": INK2}], "space_after": 0},
    ])

def add_table(slide, x, y, w, col_widths, rows, header_fill=NAVY, header_color=WHITE,
              body_size=14, header_size=13, row_h=0.5, first_col_bold=False):
    n_r, n_c = len(rows), len(rows[0])
    tbl_shape = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(row_h * n_r))
    tbl = tbl_shape.table
    tbl.first_row = False; tbl.horz_banding = False
    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = Inches(cw)
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else GRAY
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run(); run.text = val
            f = run.font
            f.size = Pt(header_size if r == 0 else body_size)
            f.bold = (r == 0) or (first_col_bold and c == 0 and r > 0)
            f.color.rgb = header_color if r == 0 else INK
            f.name = FONT
    return tbl

def arrow(slide, x, y, w=0.38, h=0.34, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

# ================================================================ SLIDE 1 — Title
s = new_slide()
bg(s)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, SH - 0.28, SW, 0.28, ORANGE)
text(s, 0.9, 1.35, 11.5, 1.6, [
    {"runs": [{"t": "Sonata Knowledge Assistant", "size": 44, "bold": True, "color": WHITE}]},
    {"runs": [{"t": "A citation-grounded answer engine for Sonata functional questions", "size": 20, "color": RGBColor(0xD5, 0xE1, 0xF3)}], "space_after": 0},
])
text(s, 0.9, 3.4, 11.5, 0.6, [
    {"runs": [{"t": "Budget proposal · Phase 1 — Functional Knowledge Base", "size": 17, "bold": True, "color": ORANGE}]},
])
# stat strip
tiles = [
    ("27/27", "eval questions correct (100%)", AQUA),
    ("100%", "answers carried accurate citations", BLUE),
    ("8", "live demo questions in ~30s", ORANGE),
    ("~$0.11", "API cost per question (measured)", GOLD),
]
tw, th = 2.72, 1.15
tx, ty = 0.9, 4.55
for i, (v, l, ac) in enumerate(tiles):
    x = tx + i * (tw + 0.18)
    rect(s, x, ty, tw, th, fill=RGBColor(0x1A, 0x52, 0x9E), round_=True, radius=0.09)
    text(s, x + 0.18, ty + 0.16, tw - 0.36, th - 0.3, [
        {"runs": [{"t": v, "size": 26, "bold": True, "color": WHITE}]},
        {"runs": [{"t": l, "size": 11, "color": RGBColor(0xCF, 0xDD, 0xF2)}], "space_after": 0},
    ])
text(s, 0.9, 6.35, 11.5, 0.5, [
    {"runs": [{"t": "Pilot complete · SME sign-off obtained · built on Bravura's own Wiki + Jira · read-only", "size": 13, "color": RGBColor(0xBF, 0xCE, 0xE8)}], "space_after": 0},
])

# ================================================================ SLIDE 2 — Problem
s = new_slide()
header(s, "The problem we are trying to solve", 2)
bullets(s, 0.7, 1.75, 12.0, 4.8, [
    "Functional questions about Sonata — \"how does X work?\", \"what changed in v16.2?\" — are answered today by manually searching Wiki + Jira, or by asking the team.",
    [{"t": "Slow and inconsistent: ", "bold": True}, {"t": "answers depend on who you ask and what they remember; knowledge lives in people's heads and scattered tickets.", }],
    [{"t": "No single trusted source: ", "bold": True}, {"t": "the same question gets different answers; undocumented behaviour is only discovered the hard way.", }],
    [{"t": "Repeat questions burn SME time: ", "bold": True}, {"t": "support, servicedesk and BAs re-ask the same things instead of self-serving.", }],
    [{"t": "Upgrade-impact questions are the hardest: ", "bold": True}, {"t": "\"which clients are affected by this change?\" takes days of manual cross-referencing — and that needs this corpus as a foundation.", }],
], size=16, gap=16)
rect(s, 0.7, 5.7, 12.0, 0.9, fill=BLUE_BG, round_=True, radius=0.12)
text(s, 0.95, 5.88, 11.5, 0.6, [
    {"runs": [{"t": "Goal: ", "bold": True, "color": NAVY}, {"t": "make every Sonata functional answer fast, consistent, and traceable to a source — then build the data layer upgrade impact needs.", "color": INK}], "space_after": 0},
])

# ================================================================ SLIDE 3 — What it is
s = new_slide()
header(s, "What Sonata Knowledge Assistant is", 3)
text(s, 0.7, 1.6, 12.0, 0.9, [
    {"runs": [{"t": "A chat assistant that answers Sonata functional questions, grounded in Bravura's own Wiki + Jira — ", },
              {"t": "never from the model's memory alone.", "bold": True}], "space_after": 0},
])
cards = [
    ("Citations on every answer", "Each answer names its source — a Wiki page/section or a Jira key — so you can open it and verify.", BLUE),
    ("Read-only by design", "Never writes back to Wiki, Jira, Bitbucket or X-ray. Outputs are drafts for human review.", AQUA),
    ("Honest when it doesn't know", "If retrieval finds nothing relevant, it says so plainly instead of guessing.", ORANGE),
    ("Scales from pilot to product", "Proven on one feature; Phase 1 extends it to all modules — the foundation for upgrade impact (Phase 3).", NAVY),
]
cw, ch = 6.0, 1.5
pos = [(0.7, 2.75), (6.83, 2.75), (0.7, 4.45), (6.83, 4.45)]
for (t, d, ac), (x, y) in zip(cards, pos):
    rect(s, x, y, cw, ch, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.07)
    rect(s, x, y, 0.14, ch, ac)
    text(s, x + 0.35, y + 0.18, cw - 0.6, ch - 0.35, [
        {"runs": [{"t": t, "size": 17, "bold": True, "color": INK}]},
        {"runs": [{"t": d, "size": 13, "color": INK2}], "space_after": 0},
    ])

# ================================================================ SLIDE 4 — How it works
s = new_slide()
header(s, "How it works — the trusted-answer loop", 4)
steps = [
    ("1 · Ingest", "Read-only pull of Wiki pages + Jira tickets; split into sections/chunks.", BLUE),
    ("2 · Index", "Chunks embedded + keyword-indexed in a retrieval store.", AQUA),
    ("3 · Retrieve", "Question → most relevant chunks found (semantic + keyword hybrid).", ORANGE),
    ("4 · Answer", "Model composes the answer from those chunks only, with inline citations.", NAVY),
]
bx, by, bw, bh = 0.62, 2.35, 2.7, 2.35
for i, (t, d, ac) in enumerate(steps):
    x = bx + i * (bw + 0.42)
    rect(s, x, by, bw, bh, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.07)
    rect(s, x, by, bw, 0.11, ac)
    text(s, x + 0.2, by + 0.3, bw - 0.4, bh - 0.5, [
        {"runs": [{"t": t, "size": 16, "bold": True, "color": INK}]},
        {"runs": [{"t": d, "size": 12.5, "color": INK2}], "space_after": 0},
    ])
    if i < 3:
        arrow(s, x + bw + 0.08, by + bh / 2 - 0.17, 0.26, 0.34, ac)
rect(s, 0.7, 5.15, 12.0, 1.05, fill=BLUE_BG, round_=True, radius=0.1)
text(s, 0.95, 5.32, 11.5, 0.75, [
    {"runs": [{"t": "The rule that makes it trustworthy: ", "bold": True, "color": NAVY},
              {"t": "if retrieval returns nothing relevant, the assistant says so — it does not fall back to general knowledge about Sonata.", "color": INK}], "space_after": 0},
])

# ================================================================ SLIDE 5 — Proof
s = new_slide()
header(s, "Proof it works — the pilot", 5)
tiles = [
    ("27/27", "eval questions correct (100%)", AQUA, None),
    ("100%", "answers carried accurate citations", BLUE, None),
    ("8", "live demo questions in ~30s", ORANGE, None),
    ("~$0.11", "API cost per question", GOLD, None),
]
tw, th = 2.72, 1.12
tx, ty = 0.7, 1.55
for i, (v, l, ac, vc) in enumerate(tiles):
    stat_tile(s, tx + i * (tw + 0.18), ty, tw, th, v, l, accent=ac, value_color=vc, value_size=30)
text(s, 0.7, 2.95, 12.0, 0.4, [
    {"runs": [{"t": "Example answers from the live demo (every one graded Correct by SME Pratigya, with the source cited):", "size": 13, "bold": True, "color": INK2}], "space_after": 0},
])
examples = [
    ("Q:  What is the default number of results per page if pagingRange is not supplied?",
     "A:  20 results per page, starting at result index 1.",
     "Wiki: RLSI-6059 § SFC-04"),
    ("Q:  Does the searchEmployer pagination change return Employer External References?",
     "A:  No — returning external references was dropped from scope (struck through in the design).",
     "Wiki: RLSI-6059 § Out of Scope"),
    ("Q:  What exception does BASE-460272 report during CreateEmployerAccount?",
     "A:  HTTP 500 SOAP fault rooted in java.lang.StackOverflowError in EmployersearchDO#retrieve().",
     "Jira: BASE-460272"),
]
ey = 3.45
for i, (q, a, c) in enumerate(examples):
    y = ey + i * 1.02
    rect(s, 0.7, y, 12.0, 0.9, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.1)
    text(s, 0.95, y + 0.1, 11.6, 0.75, [
        {"runs": [{"t": q + "   ", "size": 12, "bold": True, "color": NAVY}]},
        {"runs": [{"t": a + "  ", "size": 12.5, "color": INK}]},
        {"runs": [{"t": "  " + c, "size": 10.5, "italic": True, "color": MUTED}]},
    ])

# ================================================================ SLIDE 6 — Benefits
s = new_slide()
header(s, "Why it matters — benefits", 6)
benefits = [
    ("Faster, self-serve answers", "Support, servicedesk and BAs answer functional questions across all modules without digging through Wiki + Jira or waiting on SMEs.", BLUE),
    ("One trusted, verifiable source", "Every answer is traceable to a specific wiki section or Jira key — the trust foundation for anything client-facing later.", AQUA),
    ("Surfaces documentation debt", "Reveals which modules are well-documented and which aren't — a concrete roadmap for improving the source docs.", ORANGE),
    ("Measurable quality baseline", "Thumbs up/down feedback and an \"I don't know\" escalation rate give a tracked quality metric, not assumptions.", NAVY),
    ("Unlocks upgrade impact assessment", "Phase 3 (which clients are affected by a version change) needs this full indexed history — the single biggest downstream value.", GOOD),
]
for i, (t, d, ac) in enumerate(benefits):
    y = 1.7 + i * 1.06
    rect(s, 0.7, y, 12.0, 0.92, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.1)
    rect(s, 0.7, y, 0.14, 0.92, ac)
    text(s, 1.05, y + 0.13, 11.4, 0.7, [
        {"runs": [{"t": t + "  —  ", "size": 15, "bold": True, "color": INK},
                  {"t": d, "size": 13.5, "color": INK2}], "space_after": 0},
    ])

# ================================================================ SLIDE 7 — Architecture
s = new_slide()
header(s, "Architecture — layered, read-only, extensible", 7)
layers = [
    ("Interface", "Chat (web / Teams / Slack) · voice later (Phase 4) — thin presentation, no knowledge logic.", BLUE),
    ("Orchestration", "Intent routing (functional Q&A vs version-diff vs impact), conversation memory, citation enforcement.", AQUA),
    ("Retrieval", "Hybrid search — semantic + keyword + metadata filters — with per-source retrievers (Wiki/Jira/Bitbucket/X-ray).", ORANGE),
    ("Knowledge & index", "Chunked, embedded docs per source + entity graph: Feature ↔ Wiki ↔ Jira ↔ code ↔ test ↔ release version.", NAVY),
    ("Ingestion (read-only)", "Connectors to Confluence (Wiki) and Jira APIs — pull only, never write back.", GOLD),
]
ly = 1.65
lh = 0.92
for i, (name, desc, ac) in enumerate(layers):
    y = ly + i * (lh + 0.09)
    rect(s, 0.7, y, 12.0, lh, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.09)
    rect(s, 0.7, y, 2.5, lh, ac, round_=True, radius=0.09)
    text(s, 0.9, y + 0.24, 2.2, 0.5, [{"runs": [{"t": name, "size": 15, "bold": True, "color": WHITE}], "space_after": 0}])
    text(s, 3.45, y + 0.16, 9.1, 0.65, [{"runs": [{"t": desc, "size": 12.5, "color": INK}], "space_after": 0}], anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.7, 6.55, 12.0, 0.55, fill=BLUE_BG, round_=True, radius=0.14)
text(s, 0.95, 6.68, 11.5, 0.35, [
    {"runs": [{"t": "Read-only against all source systems — nothing is ever written back.", "size": 12.5, "bold": True, "color": NAVY}], "space_after": 0},
])

# ================================================================ SLIDE 8 — Tech stack
s = new_slide()
header(s, "Technology stack — proven, low-friction", 8)
add_table(s, 0.7, 1.6, 12.0, [3.6, 8.4], [
    ["Layer", "Choice"],
    ["LLM / orchestration", "Anthropic Claude (tool-use) with prompt caching — the POC already runs on it"],
    ["Retrieval", "Hybrid search: semantic + keyword; local index today, pluggable vector store (use an existing org license)"],
    ["Data connectors", "Confluence (Wiki) API · Jira API — read-only, citation metadata preserved"],
    ["Ingestion & chunking", "Python — wiki/Jira ingestion scripts, section-aware chunking"],
    ["Interface", "Simple internal web chat (no auth complexity for now); voice in Phase 4"],
    ["Phase 2 additions", "Bitbucket (code) + X-ray (tests) connectors — same pattern"],
], row_h=0.62, body_size=13.5)
text(s, 0.7, 5.6, 12.0, 1.2, [
    {"runs": [{"t": "No new procurement needed: ", "bold": True}, {"t": "the stack builds on APIs Bravura already has and an LLM service already in use. Deployment stays internal — no client-facing exposure in this phase.", "color": INK}], "space_after": 0},
])

# ================================================================ SLIDE 9 — Roadmap
s = new_slide()
header(s, "Roadmap — what this unlocks", 9)
phases = [
    ("Phase 0 — POC", "One feature · citations proven · SME sign-off obtained", GOOD, "DONE"),
    ("Phase 1 — Functional KB (this proposal)", "Full Wiki + Jira · all modules · 100+ question eval", BLUE, "NEXT"),
    ("Phase 2 — Code + test layer", "Bitbucket PRs + X-ray tests, linked into the graph", INK2, "later"),
    ("Phase 3 — Upgrade Impact Assistant", "Version-diff, client profiles, impact/risk reports", ORANGE, "goal"),
    ("Phase 4+ — Voice · hardening · continuous improvement", "STT/TTS interface · access control · gap detection", NAVY, "later"),
]
for i, (name, desc, ac, tag) in enumerate(phases):
    y = 1.6 + i * 0.98
    rect(s, 0.7, y, 12.0, 0.86, fill=CARD, line=HAIR, line_w=1, round_=True, radius=0.1)
    rect(s, 0.7, y, 0.14, 0.86, ac)
    rect(s, 0.98, y + 0.16, 1.15, 0.5, fill=ac, round_=True, radius=0.3)
    text(s, 0.98, y + 0.27, 1.15, 0.3, [{"runs": [{"t": tag, "size": 11, "bold": True, "color": WHITE}], "align": PP_ALIGN.CENTER, "space_after": 0}])
    text(s, 2.4, y + 0.13, 10.1, 0.6, [
        {"runs": [{"t": name, "size": 15, "bold": True, "color": INK}]},
        {"runs": [{"t": "  ·  " + desc, "size": 12.5, "color": INK2}], "space_after": 0},
    ])
text(s, 0.7, 6.6, 12.0, 0.5, [
    {"runs": [{"t": "Phases 2 and 3 build directly on the Phase 1 corpus — this proposal is the foundation for the highest-value outcome (upgrade impact).", "size": 13, "bold": True, "color": NAVY}], "space_after": 0},
])

# ================================================================ SLIDE 10 — Cost
s = new_slide()
header(s, "Cost — small and mostly people-time", 10)
stat_tile(s, 0.7, 1.6, 5.9, 1.35, "6–8 weeks", "Phase 1 effort · ~1 engineer + part-time SME/BA (per roadmap)", accent=NAVY, value_size=28)
stat_tile(s, 6.85, 1.6, 5.85, 1.35, "~$0.11", "API cost per question — measured in the pilot", accent=ORANGE, value_size=28)
add_table(s, 0.7, 3.3, 12.0, [4.4, 7.6], [
    ["Cost item", "Estimate"],
    ["API — one-time eval & tuning (100+ questions, several passes)", "~$50–150"],
    ["API — ongoing live usage (50–100 internal questions/day)", "~$100–250 / month"],
    ["API — indexing the corpus", "~$0 today (local index); optional paid vector store later"],
    ["Effort — the real cost", "~6–8 weeks: 1 engineer + part-time SME/BA for taxonomy & grading"],
], row_h=0.55, body_size=13.5)
text(s, 0.7, 5.6, 12.0, 1.0, [
    {"runs": [{"t": "Bottom line: ", "bold": True}, {"t": "the API bill stays in the low hundreds of dollars per month. The investment is people-time (engineering + SME workshops) — which also produces the documentation-debt roadmap as a side effect.", "color": INK}], "space_after": 0},
])

# ================================================================ SLIDE 11 — Ask
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.16, NAVY)
text(s, 0.7, 1.0, 12.0, 0.7, [
    {"runs": [{"t": "The ask", "size": 26, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.95, 12.0, 3.4, [
    [{"t": "Approve Phase 1 ", "bold": True}, {"t": "(6–8 weeks, ~1 engineer + SME time) to scale the proven pilot into a full Functional Knowledge Base.", }],
    [{"t": "Kick off now: ", "bold": True}, {"t": "book the Feature/Module taxonomy workshop (the critical path), enumerate the Wiki/Jira corpus, and start Jira bulk-ingestion.", }],
    [{"t": "Expected outcome: ", "bold": True}, {"t": "broad functional Q&A with citations across all modules, a measured quality baseline — and the data foundation for upgrade impact assessment.", }],
], size=16, gap=18)
rect(s, 0.7, 5.15, 12.0, 1.15, fill=BLUE_BG, round_=True, radius=0.1)
text(s, 0.95, 5.35, 11.5, 0.8, [
    {"runs": [{"t": "The pilot already delivers: ", "bold": True, "color": NAVY},
              {"t": "27/27 correct, 100% cited, SME-approved. Phase 1 just turns that proven pattern on for every module.", "color": INK}], "space_after": 0},
])
text(s, 0.7, 6.6, 12.0, 0.5, [
    {"runs": [{"t": "Thank you — questions & discussion welcome.", "size": 14, "italic": True, "color": INK2}], "space_after": 0},
])

# ---------------------------------------------------------------- save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Sonata-Knowledge-Assistant-Phase1-Proposal.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
