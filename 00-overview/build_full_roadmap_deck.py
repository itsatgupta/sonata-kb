"""Generate the Sonata Knowledge Assistant (SKA) Full Roadmap deck.

Shows all phases (0-4) with capabilities, timeline, and business value.
POLISHED VERSION: improved spacing, borders, typography, and accessibility.

Run with the POC venv python:
    "C:\\Users\\atgupta\\sonata-kb\\03-poc\\agent\\venv\\Scripts\\python.exe" build_full_roadmap_deck.py
Output: 00-overview/Sonata-Knowledge-Assistant-Full-Roadmap.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================ palette
INK      = RGBColor(0x0B, 0x0B, 0x0B)
INK2     = RGBColor(0x52, 0x51, 0x4E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PAGE     = RGBColor(0xFC, 0xFC, 0xFB)
NAVY     = RGBColor(0x10, 0x42, 0x81)
BLUE     = RGBColor(0x2A, 0x78, 0xD6)
BLUE_BG  = RGBColor(0xE4, 0xEE, 0xFB)
ORANGE   = RGBColor(0xEB, 0x68, 0x34)
AQUA     = RGBColor(0x1B, 0xAF, 0x7A)
GOLD     = RGBColor(0xED, 0xA1, 0x00)
GOOD     = RGBColor(0x0C, 0xA3, 0x0C)
GRAY     = RGBColor(0xF2, 0xF3, 0xF5)
HAIR     = RGBColor(0xE1, 0xE0, 0xD9)

FONT = "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, color=PAGE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, round_=False, border=None):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    elif line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if round_:
        shape.adjustments[0] = 0.1
    return shape

def text(s, x, y, w, h, paragraphs, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)

    for i, p_dict in enumerate(paragraphs if isinstance(paragraphs, list) else [paragraphs]):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(p_dict.get("space_after", 10))
        p.level = p_dict.get("level", 0)

        for run_dict in p_dict.get("runs", [{"t": str(p_dict)}]):
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", bold)
            r.font.color.rgb = run_dict.get("color", color)
            r.font.name = FONT
            r.font.italic = run_dict.get("italic", False)
    return box

def bullets(s, x, y, w, h, bullets_list, size=13):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)
    for i, bullet_runs in enumerate(bullets_list):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        for run_dict in bullet_runs:
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", False)
            r.font.color.rgb = run_dict.get("color", INK)
            r.font.name = FONT

# ============================================================ SLIDE 1: Title
s = new_slide()
bg(s, PAGE)
rect(s, 0, 0, SW, SH, fill=NAVY, border=HAIR)
rect(s, 0, SH - 0.28, SW, 0.28, fill=ORANGE)
text(s, 0.9, 1.2, 11.5, 2.2, [
    {"runs": [{"t": "Sonata Knowledge Assistant", "size": 52, "bold": True, "color": WHITE}]},
    {"runs": [{"t": "From POC to Production: The Full Roadmap", "size": 24, "color": RGBColor(0xD5, 0xE1, 0xF3)}], "space_after": 0},
])
text(s, 0.9, 3.6, 11.5, 0.6, [
    {"runs": [{"t": "5 phases · Citation-grounded AI · Read-only foundation", "size": 17, "color": ORANGE, "bold": True}], "space_after": 0},
])

# ============================================================ SLIDE 2: Problem & Vision
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "The Problem", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Sonata knowledge fragmented across 4 systems"}],
    [{"t": "No unified query layer — answering functional questions requires manual SME cross-referencing"}],
    [{"t": "Upgrade assessments take 2–3 days per client (vs. 1 hour with the assistant)"}],
    [{"t": "Bottleneck: same ~5 senior engineers field every functional question"}],
])

# ============================================================ SLIDE 3: Phase 0 (Complete)
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 0: POC — ", "size": 30, "bold": True, "color": NAVY},
              {"t": "✓ COMPLETE", "size": 28, "bold": True, "color": GOOD}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "What We Built", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Single feature (searchEmployer pagination)"}],
    [{"t": "20–30 curated test questions"}],
    [{"t": "Proof: ingest → index → retrieve → cite"}],
    [{"t": "Cost: ~$0.11/question"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=GRAY, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Results", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "27/27 correct (100%)"}],
    [{"t": "100% accurate citations"}],
    [{"t": "SME Pratigya: sign-off obtained ✓"}],
    [{"t": "Ready to scale to all modules"}],
], size=11.5)

# ============================================================ SLIDE 4: Phase 1
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 1: Functional Knowledge Base (Chat)", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=AQUA, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Scope", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "All Wiki spaces (design, arch, specs)", "color": WHITE}],
    [{"t": "All Jira projects (stories, defects)", "color": WHITE}],
    [{"t": "Feature/Module taxonomy (SME workshop)", "color": WHITE}],
    [{"t": "Multi-turn chat interface", "color": WHITE}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=GOLD, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Delivers", "bold": True, "size": 13, "color": INK}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Answer functional Q&A across all modules", "color": INK}],
    [{"t": "100+ test questions (≥80% correct)", "color": INK}],
    [{"t": "Feedback signal (up/down votes)", "color": INK}],
    [{"t": "Timeline: 6–8 weeks (1 eng + SME)", "color": INK}],
], size=11.5)

# ============================================================ SLIDE 5: Phase 2
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 2: Code + Test Layer (Bitbucket + X-ray)", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "New Data", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Bitbucket PR metadata (modules, merge date)"}],
    [{"t": "X-ray test cases + history"}],
    [{"t": "Entity graph: Wiki ↔ Jira ↔ PR ↔ Test"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=ORANGE, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Enables", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "\"What test coverage for feature X?\"", "color": WHITE}],
    [{"t": "Change traceability: feature → code → tests", "color": WHITE}],
    [{"t": "Data foundation for Phase 3", "color": WHITE}],
    [{"t": "Timeline: ~6 weeks", "color": WHITE}],
], size=11.5)

# ============================================================ SLIDE 6: Phase 3
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 3: Upgrade Impact Assessment", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=GOLD, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "The Ask", "bold": True, "size": 13, "color": INK}], "space_after": 0},
])
text(s, 0.95, 2.7, 5.3, 4.0, [
    {"runs": [{"t": "Given: client, from_version, to_version", "size": 11, "color": INK}]},
    {"runs": [{"t": "\n"}]},
    {"runs": [{"t": "Generate: what changed, what's relevant, what's risky", "size": 11, "color": INK}]},
    {"runs": [{"t": "\n"}]},
    {"runs": [{"t": "Output: draft report for human sign-off", "size": 11, "color": INK}]},
    {"runs": [{"t": "\n"}]},
    {"runs": [{"t": "Impact: 2–3 days → 1 hour per client", "size": 11, "bold": True, "color": INK}], "space_after": 0},
])

rect(s, 6.8, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "New Pieces", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Client profile data (version, modules used)"}],
    [{"t": "Risk/impact scoring model"}],
    [{"t": "Version-diff engine (what changed)"}],
    [{"t": "Word/PDF export for client hand-off"}],
    [{"t": "Timeline: ~8 weeks"}],
], size=11.5)

# ============================================================ SLIDE 7: Phase 4 + Timeline
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 4: Voice Interface + Full Timeline", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
text(s, 0.7, 1.85, 12.0, 0.35, [
    {"runs": [{"t": "Phase 4: Voice", "bold": True, "size": 12, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 2.3, 12.0, 1.3, [
    [{"t": "STT input + TTS output · Response formatting (spoken sentences, natural numbers) · Opt-in per response"}],
], size=11)

# Timeline
rect(s, 0.7, 3.85, 11.6, 0.1, fill=NAVY)
phases = [("P0 ✓\n3–4w", 0.7, 1.1, GOOD), ("P1\n6–8w", 1.95, 1.3, AQUA), ("P2\n6w", 3.4, 1.0, GOLD), ("P3\n8w", 4.6, 1.2, ORANGE), ("P4\n4w", 5.95, 0.8, BLUE)]
for label, x, w, color in phases:
    rect(s, x, 3.95, w, 0.45, fill=color, round_=True)
    text(s, x, 4.02, w, 0.35, [
        {"runs": [{"t": label, "bold": True, "size": 10, "color": WHITE if color in [NAVY, ORANGE] else INK}], "space_after": 0},
    ], align=PP_ALIGN.CENTER)

text(s, 0.7, 4.6, 12.0, 0.3, [
    {"runs": [{"t": "Effort per phase (indicative)", "bold": True, "size": 11, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 5.0, 12.0, 2.3, [
    [{"t": "P1: 1 eng + SME (long pole: taxonomy workshop) · P2: 1 eng · P3: 1–2 engs + domain experts · P4: 0.5 eng (parallelizable)"}],
], size=10)

# ============================================================ SLIDE 8: Business Value
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Business Value at Each Phase", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
value_rows = [
    ("Phase 1", "Support/QA/BA can self-serve answers → fewer SME escalations", AQUA, WHITE),
    ("Phase 2", "Answer technical questions about code changes & test coverage → safer support", GOLD, INK),
    ("Phase 3", "Assessment: 2–3 days → 1 hour per client → 10–15 clients/year unblocked", ORANGE, WHITE),
    ("Phase 4", "Voice option for field/mobile: hands-free Q&A on client calls", BLUE, WHITE),
]
y = 1.9
for phase, value, color, text_color in value_rows:
    rect(s, 0.7, y - 0.05, 11.6, 0.95, fill=color, round_=True)
    text(s, 0.95, y + 0.1, 2.8, 0.75, [
        {"runs": [{"t": phase, "bold": True, "size": 11, "color": text_color}], "space_after": 0}
    ])
    text(s, 3.85, y + 0.1, 8.5, 0.75, [
        {"runs": [{"t": value, "size": 11, "color": text_color}], "space_after": 0}
    ])
    y += 1.1

# ============================================================ SLIDE 9: Key Dependencies
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Key Dependencies & Risks", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.85, 12.0, 2.8, [
    [{"t": "Phase 1 → 3: Feature/Module taxonomy (SME workshop) is the ", "bold": False},
     {"t": "long pole", "bold": True}],
    [{"t": "Client profile data (Phase 3): doesn't exist cleanly — requires interviews or config review"}],
    [{"t": "Wiki/Jira quality uneven — Phase 1 surfaces documentation debt"}],
    [{"t": "Version-to-trunk mapping: must validate fixVersion accuracy"}],
], size=12)
rect(s, 0.7, 4.85, 11.6, 2.4, fill=BLUE_BG, round_=True)
text(s, 0.95, 5.05, 11.2, 0.4, [
    {"runs": [{"t": "Mitigation", "bold": True, "size": 12, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.95, 5.55, 11.2, 1.6, [
    [{"t": "Make escalation / \"I don't know\" cheap & explicit (not a failure — a feature)"}],
    [{"t": "Surface last-updated dates on every citation (flag stale content)"}],
    [{"t": "Documentation-debt backlog is valid Phase 1 output (drives wiki improvements)"}],
], size=11)

# ============================================================ SLIDE 10: Next Steps
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Next Steps: Phase 1 Kickoff", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.9, 12.0, 2.8, [
    [{"t": "Book the Feature/Module taxonomy workshop (the critical path item)"}],
    [{"t": "Enumerate all Wiki spaces + Jira projects for full ingestion scope"}],
    [{"t": "Build Jira bulk-ingestion script"}],
    [{"t": "Apply cost-optimization levers (validated on Phase 0)"}],
], size=13)
rect(s, 0.7, 5.0, 11.6, 2.2, fill=ORANGE, round_=True)
text(s, 0.95, 5.2, 11.2, 0.4, [
    {"runs": [{"t": "The Opportunity", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
text(s, 0.95, 5.75, 11.2, 1.3, [
    {"runs": [{"t": "Phase 0 proved it works: 27/27 correct, 100% cited. Phase 1 just turns that proven pattern on for every module. From there, Phases 2–4 unlock client-facing upgrade assessment.", "size": 12, "color": WHITE}], "space_after": 0},
])

# ============================================================ save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Sonata-Knowledge-Assistant-Full-Roadmap.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides (polished)")
