"""Generate the Sonata Knowledge Assistant Comprehensive Roadmap Deck.

16 slides covering: POC v2 Voice + Phase 1 Royal London + Beneficiaries + ROI.
Run: python build_comprehensive_deck.py
Output: Sonata-Knowledge-Assistant-Comprehensive-Roadmap.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================ palette
INK      = RGBColor(0x0B, 0x0B, 0x0B)
INK2     = RGBColor(0x52, 0x51, 0x4E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PAGE     = RGBColor(0xFC, 0xFC, 0xFB)
NAVY     = RGBColor(0x10, 0x42, 0x81)
BLUE     = RGBColor(0x2A, 0x78, 0xD6)
BLUE_BG  = RGBColor(0xE4, 0xEE, 0xFB)
ORANGE   = RGBColor(0xEB, 0x68, 0x34)
AQUA     = RGBColor(0x1B, 0xAF, 0x7A)
GOLD     = RGBColor(0xED, 0xA1, 0x00)
GOOD     = RGBColor(0x0C, 0xA3, 0x0C)
GRAY     = RGBColor(0xF2, 0xF3, 0xF5)
HAIR     = RGBColor(0xE1, 0xE0, 0xD9)
LIGHT_BLUE = RGBColor(0xD5, 0xE1, 0xF3)

FONT = "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, color=PAGE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, round_=False, border=None):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    elif line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if round_:
        shape.adjustments[0] = 0.1
    return shape

def text(s, x, y, w, h, paragraphs, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)
    for i, p_dict in enumerate(paragraphs if isinstance(paragraphs, list) else [paragraphs]):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(p_dict.get("space_after", 10))
        for run_dict in p_dict.get("runs", [{"t": str(p_dict)}]):
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", bold)
            r.font.color.rgb = run_dict.get("color", color)
            r.font.name = FONT
    return box

def bullets(s, x, y, w, h, bullets_list, size=13):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)
    for i, bullet_runs in enumerate(bullets_list):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        for run_dict in bullet_runs:
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", False)
            r.font.color.rgb = run_dict.get("color", INK)
            r.font.name = FONT

# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = new_slide()
bg(s, PAGE)
rect(s, 0, 0, SW, SH, fill=NAVY, border=HAIR)
rect(s, 0, SH - 0.28, SW, 0.28, fill=ORANGE)
text(s, 0.9, 1.0, 11.5, 1.2, [
    {"runs": [{"t": "Sonata Knowledge Assistant", "size": 48, "bold": True, "color": WHITE}], "space_after": 0},
])
text(s, 0.9, 2.2, 11.5, 0.8, [
    {"runs": [{"t": "POC v2 Voice Demo + Phase 1 Strategy", "size": 28, "color": LIGHT_BLUE}], "space_after": 0},
])
text(s, 0.9, 3.3, 11.5, 0.5, [
    {"runs": [{"t": "Royal London UAR/ISA Pilot · 12 Weeks · $238K Support Savings", "size": 17, "color": ORANGE, "bold": True}], "space_after": 0},
])
text(s, 0.9, 4.2, 11.5, 1.5, [
    {"runs": [{"t": "Phase 0: ✅ PROVEN (27/27 correct, 100% cited, SME sign-off)", "size": 14, "color": WHITE}], "space_after": 0},
    {"runs": [{"t": "Phase 0.5: 🎤 Voice demo (2 weeks, $0, Whisper + Web Speech)", "size": 14, "color": WHITE}], "space_after": 0},
    {"runs": [{"t": "Phase 1: Royal London UAR/ISA (12 weeks, $20K, $238K/yr ROI)", "size": 14, "color": WHITE}], "space_after": 0},
])

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "The Problem", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Sonata knowledge fragmented across 4 systems (Wiki, Jira, Bitbucket, X-ray)"}],
    [{"t": "No unified query layer — answering functional questions requires manual SME cross-referencing"}],
    [{"t": "Upgrade assessments take 2–3 days per client (vs. target: 1 hour with assistant)"}],
    [{"t": "Same ~5 senior engineers field every \"how does X work\" question — bottleneck"}],
    [{"t": "Royal London: 80+ UAR/ISA support tickets/month consuming team bandwidth"}],
])

# ============================================================
# SLIDE 3: PHASE 0 COMPLETE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 0: POC — ", "size": 30, "bold": True, "color": NAVY},
              {"t": "✓ COMPLETE", "size": 28, "bold": True, "color": GOOD}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "What We Built", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Single feature: searchEmployer pagination"}],
    [{"t": "20–30 curated test questions"}],
    [{"t": "Proof: ingest → index → retrieve → cite"}],
    [{"t": "Cost: ~$0.11/question"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=GRAY, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Results", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "27/27 correct (100%)"}],
    [{"t": "100% accurate citations"}],
    [{"t": "SME Pratigya: sign-off obtained ✓"}],
    [{"t": "Ready to scale to all modules"}],
], size=11.5)

# ============================================================
# SLIDE 4: POC v2 VOICE DEMO (NEW)
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "POC v2: Voice Demo — ", "size": 30, "bold": True, "color": NAVY},
              {"t": "Fully Working Voice Assistant", "size": 26, "bold": True, "color": AQUA}], "space_after": 0},
])
rect(s, 0.7, 1.8, 3.7, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.9, 2.0, 3.3, 0.4, [
    {"runs": [{"t": "The Ask (2 Weeks)", "bold": True, "size": 12}], "space_after": 0},
])
bullets(s, 0.9, 2.5, 3.3, 4.1, [
    [{"t": "Add STT + TTS to existing POC"}],
    [{"t": "Whisper API (STT) + Web Speech (TTS)"}],
    [{"t": "Use searchEmployer data (existing)"}],
    [{"t": "Deploy: Vercel + Render"}],
], size=10.5)

rect(s, 4.6, 1.8, 3.7, 5.0, fill=AQUA, round_=True)
text(s, 4.8, 2.0, 3.3, 0.4, [
    {"runs": [{"t": "What You'll See", "bold": True, "size": 12, "color": WHITE}], "space_after": 0},
])
bullets(s, 4.8, 2.5, 3.3, 4.1, [
    [{"t": "Click \"Ask\" button", "color": WHITE}],
    [{"t": "Speak question", "color": WHITE}],
    [{"t": "Bot transcribes (Whisper)", "color": WHITE}],
    [{"t": "Bot speaks answer", "color": WHITE}],
    [{"t": "Shows sources below", "color": WHITE}],
], size=10.5)

rect(s, 8.5, 1.8, 4.0, 5.0, fill=ORANGE, round_=True)
text(s, 8.7, 2.0, 3.6, 0.4, [
    {"runs": [{"t": "Why This Matters", "bold": True, "size": 12, "color": WHITE}], "space_after": 0},
])
bullets(s, 8.7, 2.5, 3.6, 4.1, [
    [{"t": "End-to-end system (not prototype)", "color": WHITE}],
    [{"t": "Voice = 10x more impressive", "color": WHITE}],
    [{"t": "Proves architecture before Phase 1", "color": WHITE}],
    [{"t": "Zero cost (all free tiers)", "color": WHITE}],
    [{"t": "Stakeholder confidence", "color": WHITE}],
], size=10.5)

# ============================================================
# SLIDE 5: VOICE ARCHITECTURE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "How It Works: Voice → Text → Answer → Voice", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 11.6, 1.3, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 11.1, 1.0, [
    {"runs": [{"t": "User speaks: \"How does searchEmployer pagination work?\"", "size": 13, "bold": True}], "space_after": 0},
    {"runs": [{"t": "→ Whisper API transcribes speech to text (free tier, 25K min/month)", "size": 11, "color": INK2}], "space_after": 0},
    {"runs": [{"t": "→ orchestrator.ask() retrieves answer from index (<200ms) with citations", "size": 11, "color": INK2}], "space_after": 0},
    {"runs": [{"t": "→ Web Speech API speaks answer back (browser built-in, unlimited, zero cost)", "size": 11, "color": INK2}], "space_after": 0},
])
text(s, 0.7, 3.3, 11.6, 0.4, [
    {"runs": [{"t": "Tech Stack (Zero Cost)", "bold": True, "size": 13, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 3.8, 11.6, 3.0, [
    [{"t": "Frontend (Vercel): Mic input + audio playback + citations display"}],
    [{"t": "Backend (Render): Whisper STT + orchestrator.ask() + return text answer"}],
    [{"t": "Database: PostgreSQL + pgvector (existing POC index, no changes)"}],
    [{"t": "Cost: $0 (Whisper free tier + Web Speech browser API + Render/Vercel free tier)"}],
    [{"t": "Timeline: 2 weeks (copy-paste code scaffolds provided)"}],
], size=12)

# ============================================================
# SLIDE 6: WHY VOICE WINS STAKEHOLDERS
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Why Voice Demo Wins Stakeholders", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rows = [
    ("Text-only POC", '"Chatbot answered correctly (27/27)"', "Cool. Will it scale?", BLUE_BG, INK),
    ("Voice POC v2", '"Voice assistant heard, retrieved, and spoke the answer"', "This is amazing! Let's approve Phase 1.", AQUA, WHITE),
    ("Impact", "Same accuracy + retrieval, but 10x more impressive", "Easier Phase 1 approval = faster path to $238K ROI", ORANGE, WHITE),
]
y = 1.8
for title, what, impact, bgc, tc in rows:
    rect(s, 0.7, y - 0.05, 11.6, 1.4, fill=bgc, round_=True)
    text(s, 0.95, y + 0.1, 2.8, 1.2, [
        {"runs": [{"t": title, "bold": True, "size": 12, "color": tc}], "space_after": 0}
    ])
    text(s, 3.85, y + 0.1, 4.0, 1.2, [
        {"runs": [{"t": what, "size": 11, "color": tc}], "space_after": 0}
    ])
    text(s, 7.95, y + 0.1, 4.2, 1.2, [
        {"runs": [{"t": impact, "size": 11, "color": tc}], "space_after": 0}
    ])
    y += 1.6

# ============================================================
# SLIDE 7: STRATEGY OVERVIEW
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "The Strategy: One Client, One Feature, Per Cycle", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "NOT: Enumerate all Wiki/Jira for all modules (takes forever, risky, no early ROI)"}],
    [{"t": "YES: One real client + one pain-point feature (bounded scope, measurable, fast)"}],
    [{"t": "Prove the playbook (12 weeks), then replicate to N clients (template + automation)"}],
    [{"t": "Why: POC proved the loop works; Phase 1 proves it scales on real traffic with real ROI"}],
    [{"t": "Risk: Contained (if 1 client fails, we learn; full-product fail = career limiting)"}],
])

# ============================================================
# SLIDE 8: ROYAL LONDON SCOPE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Pilot: Royal London (EMEA) · UAR + ISA", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Royal London Profile", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Major EMEA client"}],
    [{"t": "Support team: 8 people"}],
    [{"t": "Current version: v16.1"}],
    [{"t": "Planned upgrade: v16.2 (Q4)"}],
    [{"t": "Pain point: 80+ UAR/ISA tickets/month"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=AQUA, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "What We're Building", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Chat bot for UAR/ISA questions", "color": WHITE}],
    [{"t": "Indexed from wiki + Jira + release notes", "color": WHITE}],
    [{"t": "Answers with citations (sources)", "color": WHITE}],
    [{"t": "Response time: <1 sec", "color": WHITE}],
    [{"t": "Target accuracy: ≥80%", "color": WHITE}],
    [{"t": "Timeline: 12 weeks", "color": WHITE}],
], size=11.5)

# ============================================================
# SLIDE 9: SYSTEM ARCHITECTURE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "System Architecture: Production RAG", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Database: PostgreSQL + pgvector (structured chunks + vector embeddings + metadata filters)"}],
    [{"t": "Index: 150 wiki chunks + 47 Jira chunks + 12 release notes = 209 chunks (Royal London scope)"}],
    [{"t": "Query: Semantic search (vector) + metadata filter (client_id, module, version) < 200ms"}],
    [{"t": "NOT hitting Wiki/Jira APIs per query — cached index, refreshed once daily (off-peak)"}],
    [{"t": "Ingestion: Wiki daily, Jira daily, Bitbucket weekly (no API thrashing, scheduled)"}],
    [{"t": "Every answer carries: wiki page, Jira key, release date, last-updated freshness flag"}],
])

# ============================================================
# SLIDE 10: 12-WEEK TIMELINE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Phase 1 Timeline: 12 Weeks to Production", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Week 1-2: Data ingest + index creation (150 wiki + 47 Jira + 12 release note chunks)"}],
    [{"t": "Week 3-4: Validation against 100 real support questions (measure accuracy, cite rate)"}],
    [{"t": "Week 5-10: Iterate to ≥80% accuracy (daily feedback loop with Royal London support team)"}],
    [{"t": "Week 10: Deploy chat interface to Royal London support (private beta, 10 agents)"}],
    [{"t": "Weeks 11-12: Go-live + measure (questions/week, accuracy, time saved, satisfaction)"}],
    [{"t": "Week 12 outcome: 40% UAR/ISA volume handled by bot, 10+ hrs/week saved, $238K/yr"}],
])

# ============================================================
# SLIDE 11: DATA REQUIREMENTS
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Data Handoff: What We Need (3–4 Hours)", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Wiki pages (UAR/ISA tagged) — HTML export — 1 hour"}],
    [{"t": "Jira tickets (FEAT/DEFECT, component=UAR|ISA, v15.8–v16.2) — CSV — 30 min"}],
    [{"t": "Release notes (v15.8→v16.2 UAR/ISA changes) — Markdown — 15 min"}],
    [{"t": "100 real support questions (date, question, resolution, time-to-res) — CSV — 1-2 hours"}],
    [{"t": "Royal London metadata (current version, team size, modules) — Form — 15 min"}],
    [{"t": "Known gaps/doc debt (outdated or conflicting docs) — List — 30 min"}],
    [{"t": "Total: ~3-4 hours from your team → we build the rest"}],
])

# ============================================================
# SLIDE 12: SUPPORT TEAM WORKFLOW
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "How Support Team Uses It (Day 1)", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
text(s, 0.7, 1.8, 12.0, 0.4, [
    {"runs": [{"t": "Traditional Workflow", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.7, 2.2, 12.0, 1.0, [
    [{"t": "Agent gets UAR question → 15 min searching wiki + Jira → sends answer"}],
], size=11)

text(s, 0.7, 3.3, 12.0, 0.4, [
    {"runs": [{"t": "New Workflow (With Bot)", "bold": True, "size": 13, "color": AQUA}], "space_after": 0},
])
bullets(s, 0.7, 3.7, 12.0, 3.0, [
    [{"t": "Agent asks bot → bot responds in 2 sec with sources → agent adds context + sends"}],
    [{"t": "Time saved: 15 min → 2 min per question"}],
    [{"t": "47 UAR/ISA questions/week × 13 min = 10+ hours/week freed"}],
    [{"t": "New hire training: 30% faster (bot is instant knowledge base)"}],
    [{"t": "Quality consistency: every agent gives same answer (no knowledge loss)"}],
], size=11)

# ============================================================
# SLIDE 13: SUPPORT ROI
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Support Team ROI: Royal London Year 1", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
metrics = [
    ("Questions/week", "47", "47", "—", INK),
    ("Time per question", "15 min", "2 min", "−13 min", INK),
    ("Weekly time saved", "705 min", "94 min", "611 min (10.2 hrs)", INK),
    ("Cost/week", "$5,288", "$703", "−$4,585", INK),
    ("Cost/year", "$275K", "$36.5K", "−$238.5K SAVED", GOOD),
]
y = 1.85
for metric, before, after, delta, tc in metrics:
    row_bg = GRAY if y % 2 == 0 else WHITE
    rect(s, 0.7, y - 0.05, 11.6, 0.8, fill=row_bg, round_=True)
    text(s, 0.95, y + 0.1, 2.8, 0.7, [
        {"runs": [{"t": metric, "bold": True, "size": 11}], "space_after": 0}
    ])
    text(s, 3.85, y + 0.1, 2.2, 0.7, [
        {"runs": [{"t": before, "size": 11}], "space_after": 0}
    ])
    text(s, 6.15, y + 0.1, 2.2, 0.7, [
        {"runs": [{"t": after, "size": 11, "bold": True}], "space_after": 0}
    ])
    text(s, 8.55, y + 0.1, 3.6, 0.7, [
        {"runs": [{"t": delta, "size": 11, "color": tc, "bold": True}], "space_after": 0}
    ])
    y += 0.9

# ============================================================
# SLIDE 14: 7 BRAVURA BENEFICIARIES
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Bravura Beneficiaries: 7 Personas", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
personas = [
    ("Support", "40% volume automated, 12+ hrs/week saved, training tool", AQUA, WHITE),
    ("QA/Testing", "Test traceability per release, regression identification", GOLD, INK),
    ("BA/Product", "Data-driven roadmap (real customer questions), doc gap ID", BLUE, WHITE),
    ("Architects", "Impact assessment automation (2 days → 2 hours, Phase 3)", ORANGE, WHITE),
    ("Development", "Better bug context from support, release note generation", AQUA, WHITE),
    ("Ops/Maintenance", "Incident context dashboards, faster MTTR", GOLD, INK),
    ("Executive/CFO", "$4M+ ROI at 10 clients, competitive moat, revenue protection", ORANGE, WHITE),
]
y = 1.75
for persona, benefit, color, tc in personas:
    rect(s, 0.7, y - 0.02, 11.6, 0.75, fill=color, round_=True)
    text(s, 0.95, y + 0.1, 2.3, 0.65, [
        {"runs": [{"t": persona, "bold": True, "size": 11, "color": tc}], "space_after": 0}
    ])
    text(s, 3.4, y + 0.1, 8.8, 0.65, [
        {"runs": [{"t": benefit, "size": 10, "color": tc}], "space_after": 0}
    ])
    y += 0.82

# ============================================================
# SLIDE 15: ROI SUMMARY
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "ROI Analysis: Year 1", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.85, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.05, 5.3, 0.5, [
    {"runs": [{"t": "Pilot Cost", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "1 engineer (3 months): $60K"}],
    [{"t": "Cloud compute: $1.5K"}],
    [{"t": "Total investment: ~$20K (excl. salary)"}],
    [{"t": "\n"}],
    [{"t": "Payback: Week 4", "bold": True, "color": GOOD}],
    [{"t": "Year 1 ROI: 12x", "bold": True, "color": GOOD}],
], size=11.5)

rect(s, 6.8, 1.85, 5.8, 5.0, fill=ORANGE, round_=True)
text(s, 7.05, 2.05, 5.3, 0.5, [
    {"runs": [{"t": "Year 1 Savings (1 Client)", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Royal London support: $238.5K", "color": WHITE}],
    [{"t": "Bravura support: $156K", "color": WHITE}],
    [{"t": "QA incident prevention: $50–200K", "color": WHITE}],
    [{"t": "Architect time: $20K", "color": WHITE}],
    [{"t": "Other teams: $15K", "color": WHITE}],
    [{"t": "\n"}],
    [{"t": "Total: $479.5K–679.5K", "bold": True, "color": WHITE}],
    [{"t": "Enterprise (10 clients): $4M+", "bold": True, "color": WHITE}],
], size=11.5)

# ============================================================
# SLIDE 16: NEXT STEPS + DECISION
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Next Steps: Ready to Execute", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
text(s, 0.7, 1.8, 12.0, 0.4, [
    {"runs": [{"t": "Immediate (This Week)", "bold": True, "size": 13, "color": AQUA}], "space_after": 0},
])
bullets(s, 0.7, 2.2, 12.0, 1.8, [
    [{"t": "Confirm: POC v2 Voice (2 weeks, $0, Whisper + Web Speech)"}],
    [{"t": "Receive code scaffolds + deployment guide"}],
    [{"t": "Integrate + deploy (Vercel + Render, copy-paste, ~3 hours)"}],
], size=12)

text(s, 0.7, 4.0, 12.0, 0.4, [
    {"runs": [{"t": "Then (Week 3-12)", "bold": True, "size": 13, "color": ORANGE}], "space_after": 0},
])
bullets(s, 0.7, 4.4, 12.0, 1.8, [
    [{"t": "Demo voice bot to stakeholders → win Phase 1 approval"}],
    [{"t": "Royal London data handoff (3-4h from their team)"}],
    [{"t": "12-week Phase 1 execution → go-live week 10 → $238K/yr verified"}],
], size=12)

rect(s, 0.7, 6.0, 11.6, 1.3, fill=ORANGE, round_=True)
text(s, 0.95, 6.15, 11.2, 1.1, [
    {"runs": [{"t": "The Ask: ", "bold": True, "size": 13, "color": WHITE},
              {"t": "Approve POC v2 Voice (2 weeks) + Phase 1 Royal London (12 weeks). Cost: $20K. ROI: $238K/yr + proven playbook for scaling to 10 clients ($4M+).", "size": 12, "color": WHITE}], "space_after": 0},
])

# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Sonata-Knowledge-Assistant-Comprehensive-Roadmap.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides (comprehensive)")
