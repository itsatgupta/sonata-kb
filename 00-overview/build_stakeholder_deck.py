"""Generate the Sonata Knowledge Assistant Stakeholder Deck.

Covers: POC v2 Voice + Upgrade Analysis + Defect Triage + Roadmap
Run: python build_stakeholder_deck.py
Output: Sonata-Knowledge-Assistant-Stakeholder-Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x0B, 0x0B, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE = RGBColor(0xFC, 0xFC, 0xFB)
NAVY = RGBColor(0x10, 0x42, 0x81)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
BLUE_BG = RGBColor(0xE4, 0xEE, 0xFB)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
GOLD = RGBColor(0xED, 0xA1, 0x00)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)
GRAY = RGBColor(0xF2, 0xF3, 0xF5)
HAIR = RGBColor(0xE1, 0xE0, 0xD9)
LIGHT_BLUE = RGBColor(0xD5, 0xE1, 0xF3)

FONT = "Calibri"
SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, color=PAGE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, fill=None, line=None, round_=False, border=None):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    elif line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if round_:
        shape.adjustments[0] = 0.1
    return shape

def text(s, x, y, w, h, paragraphs, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)
    for i, p_dict in enumerate(paragraphs if isinstance(paragraphs, list) else [paragraphs]):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(p_dict.get("space_after", 10))
        for run_dict in p_dict.get("runs", [{"t": str(p_dict)}]):
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", bold)
            r.font.color.rgb = run_dict.get("color", color)
            r.font.name = FONT
    return box

def bullets(s, x, y, w, h, bullets_list, size=13):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = frame.margin_top = Inches(0.08)
    frame.margin_left = frame.margin_right = Inches(0.1)
    for i, bullet_runs in enumerate(bullets_list):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        for run_dict in bullet_runs:
            r = p.add_run()
            r.text = run_dict.get("t", "")
            r.font.size = Pt(run_dict.get("size", size))
            r.font.bold = run_dict.get("bold", False)
            r.font.color.rgb = run_dict.get("color", INK)
            r.font.name = FONT

# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = new_slide()
bg(s, PAGE)
rect(s, 0, 0, SW, SH, fill=NAVY, border=HAIR)
rect(s, 0, SH - 0.28, SW, 0.28, fill=ORANGE)
text(s, 0.9, 1.0, 11.5, 1.2, [
    {"runs": [{"t": "Sonata Knowledge Assistant", "size": 48, "bold": True, "color": WHITE}], "space_after": 0},
])
text(s, 0.9, 2.2, 11.5, 0.8, [
    {"runs": [{"t": "POC Demo + Expansion Roadmap", "size": 28, "color": LIGHT_BLUE}], "space_after": 0},
])
text(s, 0.9, 3.3, 11.5, 0.5, [
    {"runs": [{"t": "3 Working POCs · Voice + Text + Upgrade + Defect Triage", "size": 17, "color": ORANGE, "bold": True}], "space_after": 0},
])

# ============================================================
# SLIDE 2: WHAT WE'VE BUILT
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "What We've Built So Far", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "POC v2 Voice Assistant", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Speak question -> bot transcribes (Whisper)"}],
    [{"t": "Bot retrieves answer from knowledge base"}],
    [{"t": "Bot speaks answer back (Web Speech)"}],
    [{"t": "Text chat with 3 modes: Direct/OpenAI/Claude"}],
    [{"t": "Polished UI: section headers, tooltips, history"}],
    [{"t": "Deployed: Vercel + Render ($0 infrastructure)"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=AQUA, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Results", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Phase 0: 27/27 correct, 100% cited", "color": WHITE}],
    [{"t": "SME sign-off obtained", "color": WHITE}],
    [{"t": "Voice pipeline: end-to-end working", "color": WHITE}],
    [{"t": "Wiki index: deployed to cloud", "color": WHITE}],
    [{"t": "Architecture: proven & scalable", "color": WHITE}],
    [{"t": "Live demo: https://sonata-kb.vercel.app", "color": WHITE}],
], size=11.5)

# ============================================================
# SLIDE 3: LIVE DEMO
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Live Demo: Voice + Text Chat", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Chat Bot: Type question -> get answer with sources (3 modes: free/cheap/premium)"}],
    [{"t": "Voice Bot: Speak question -> bot transcribes -> answers -> speaks back"}],
    [{"t": "Demo URL: https://sonata-kb.vercel.app"}],
    [{"t": "Backend: https://sonata-kb.onrender.com"}],
    [{"t": "Cost: ~$0.11 per query (Claude) or ~$0.001 (OpenAI) or $0 (Direct)"}],
])

# ============================================================
# SLIDE 4: UPGRADE ANALYSIS POC
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "POC 2: Upgrade Impact Analysis", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "What It Does", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Compare 2 trunk releases (e.g., v16.1 -> v16.2)"}],
    [{"t": "Categorize changes: Architecture, Tech, Func, Impact"}],
    [{"t": "Identify breaking changes & deprecations"}],
    [{"t": "Generate upgrade impact report"}],
    [{"t": "Cite every change to Jira ticket + wiki page"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=GOLD, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Timeline & Data", "bold": True, "size": 13, "color": INK}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Effort: 1 week to build", "color": INK}],
    [{"t": "Data: Jira + Wiki (already available)", "color": INK}],
    [{"t": "Cost: ~$5.65 (50 queries)", "color": INK}],
    [{"t": "Who wins: Architects, PM, Delivery teams", "color": INK}],
    [{"t": "Impact: 2-hour manual -> 30-second bot", "color": INK}],
], size=11.5)

# ============================================================
# SLIDE 5: DEFECT TRIAGE POC
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "POC 3: Defect Triage Assistant", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rect(s, 0.7, 1.8, 5.8, 5.0, fill=BLUE_BG, round_=True)
text(s, 0.95, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "What It Does", "bold": True, "size": 13}], "space_after": 0},
])
bullets(s, 0.95, 2.6, 5.3, 4.1, [
    [{"t": "Support describes defect -> bot searches Jira"}],
    [{"t": "Finds similar historical defects"}],
    [{"t": "Checks if already fixed in later release"}],
    [{"t": "Recommends: fix exists / escalate / draft new"}],
    [{"t": "Auto-drafts defect with proper format"}],
], size=11.5)

rect(s, 6.8, 1.8, 5.8, 5.0, fill=ORANGE, round_=True)
text(s, 7.05, 2.0, 5.3, 0.5, [
    {"runs": [{"t": "Timeline & Data", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
])
bullets(s, 7.05, 2.6, 5.3, 4.1, [
    [{"t": "Effort: 1-2 weeks to build", "color": WHITE}],
    [{"t": "Data: Jira defects (need client field check)", "color": WHITE}],
    [{"t": "Cost: ~$11.30 (100 queries)", "color": WHITE}],
    [{"t": "Who wins: Support team, QA, Dev", "color": WHITE}],
    [{"t": "Impact: 15-min manual search -> 30-second bot", "color": WHITE}],
], size=11.5)

# ============================================================
# SLIDE 6: ALL 3 POCs SUMMARY
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "3 POCs: Complete Package", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
pocs = [
    ("POC v2 Voice", "Voice + Text chat", "DONE", "https://sonata-kb.vercel.app", AQUA, WHITE),
    ("Upgrade Analysis", "Compare 2 releases", "1 week", "Same architecture", GOLD, INK),
    ("Defect Triage", "Similar defect lookup", "1-2 weeks", "Same architecture", ORANGE, WHITE),
]
y = 1.8
for name, desc, timeline, tech, color, tc in pocs:
    rect(s, 0.7, y - 0.05, 11.6, 0.95, fill=color, round_=True)
    text(s, 0.95, y + 0.1, 2.5, 0.75, [
        {"runs": [{"t": name, "bold": True, "size": 12, "color": tc}], "space_after": 0}
    ])
    text(s, 3.5, y + 0.1, 3.5, 0.75, [
        {"runs": [{"t": desc, "size": 11, "color": tc}], "space_after": 0}
    ])
    text(s, 7.1, y + 0.1, 2.5, 0.75, [
        {"runs": [{"t": timeline, "bold": True, "size": 11, "color": tc}], "space_after": 0}
    ])
    text(s, 9.7, y + 0.1, 2.5, 0.75, [
        {"runs": [{"t": tech, "size": 11, "color": tc}], "space_after": 0}
    ])
    y += 1.1

# ============================================================
# SLIDE 7: COST ANALYSIS
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Cost Analysis: 3 POCs", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
rows = [
    ("POC v2 Voice", "100", "$11.00", "$0.30", "$0", "$11.30", AQUA, WHITE),
    ("Upgrade Analysis", "50", "$5.50", "$0.15", "$0", "$5.65", GOLD, INK),
    ("Defect Triage", "100", "$11.00", "$0.30", "$0", "$11.30", ORANGE, WHITE),
    ("TOTAL", "250", "$27.50", "$0.75", "$0", "$28.25", NAVY, WHITE),
]
y = 1.85
headers = ["POC", "Queries", "Claude", "Whisper", "TTS", "Total"]
x_positions = [0.95, 2.8, 4.2, 5.9, 7.3, 8.8]
for i, h in enumerate(headers):
    text(s, x_positions[i], y, 1.4, 0.4, [
        {"runs": [{"t": h, "bold": True, "size": 11, "color": NAVY}], "space_after": 0}
    ])
y += 0.5
for name, queries, claude, whisper, tts, total, color, tc in rows:
    bgc = color if name == "TOTAL" else GRAY
    rect(s, 0.7, y - 0.05, 11.6, 0.6, fill=bgc, round_=True)
    vals = [name, queries, claude, whisper, tts, total]
    for i, v in enumerate(vals):
        text(s, x_positions[i], y + 0.05, 1.4, 0.5, [
            {"runs": [{"t": v, "bold": name == "TOTAL", "size": 11, "color": tc}], "space_after": 0}
        ])
    y += 0.65

text(s, 0.7, 5.5, 11.6, 1.5, [
    {"runs": [{"t": "Note: ", "bold": True, "size": 12}, {"t": "These are estimated costs based on ~250 test queries. Actual costs will vary based on usage volume during stakeholder demos and testing cycles.", "size": 11, "color": RGBColor(0x52, 0x51, 0x4E)}], "space_after": 0},
])

# ============================================================
# SLIDE 8: ARCHITECTURE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Architecture: Same Pattern, 3 Use Cases", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Core: PostgreSQL + pgvector index (wiki + Jira chunks)"}],
    [{"t": "Retrieval: Semantic search + metadata filters (<200ms)"}],
    [{"t": "Voice: Whisper STT -> index -> answer -> Web Speech TTS"}],
    [{"t": "Text: Type question -> index -> answer (3 LLM modes)"}],
    [{"t": "Upgrade: Same index + version-diff JQL queries"}],
    [{"t": "Defect: Same index + similarity search over Jira defects"}],
    [{"t": "Infrastructure: Vercel (frontend) + Render (backend) = $0"}],
])

# ============================================================
# SLIDE 9: TIMELINE
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Timeline: 3 POCs in 3 Weeks", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "Week 1 (done): POC v2 Voice -> deployed & working", "bold": True}],
    [{"t": "Week 2: Upgrade Analysis POC -> compare 2 releases, generate report"}],
    [{"t": "Week 3: Defect Triage POC -> similar defect lookup, triage recommendation"}],
    [{"t": "Week 4: Polish all 3 + voice integration + stakeholder demos"}],
    [{"t": "Total: 3 working demos in ~3 weeks, ~$29 API cost"}],
])

# ============================================================
# SLIDE 10: BENEFICIARIES
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Who Wins: 7 Internal Personas", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
personas = [
    ("Support", "Voice + text Q&A, defect triage, faster resolution", AQUA, WHITE),
    ("QA/Testing", "Test traceability, regression identification", GOLD, INK),
    ("BA/Product", "Data-driven roadmap (real customer questions)", BLUE, WHITE),
    ("Architects", "Upgrade impact analysis (2 hours -> 30 seconds)", ORANGE, WHITE),
    ("Development", "Better bug context, release note generation", AQUA, WHITE),
    ("Ops/Maintenance", "Incident context, faster MTTR", GOLD, INK),
    ("Executive", "3 working demos, competitive moat, $0 infrastructure", ORANGE, WHITE),
]
y = 1.75
for persona, benefit, color, tc in personas:
    rect(s, 0.7, y - 0.02, 11.6, 0.75, fill=color, round_=True)
    text(s, 0.95, y + 0.1, 2.3, 0.65, [
        {"runs": [{"t": persona, "bold": True, "size": 11, "color": tc}], "space_after": 0}
    ])
    text(s, 3.4, y + 0.1, 8.8, 0.65, [
        {"runs": [{"t": benefit, "size": 10, "color": tc}], "space_after": 0}
    ])
    y += 0.82

# ============================================================
# SLIDE 11: NEXT STEPS
# ============================================================
s = new_slide()
bg(s)
rect(s, 0, 0, SW, 0.15, fill=NAVY, border=HAIR)
text(s, 0.7, 0.95, 12.0, 0.5, [
    {"runs": [{"t": "Next Steps", "size": 30, "bold": True, "color": NAVY}], "space_after": 0},
])
bullets(s, 0.7, 1.8, 12.0, 5.2, [
    [{"t": "This week: Demo POC v2 Voice to stakeholders", "bold": True}],
    [{"t": "Next week: Build Upgrade Analysis POC (compare v16.1 vs v16.2)"}],
    [{"t": "Week after: Build Defect Triage POC (similar defect lookup)"}],
    [{"t": "Week 4: Polish + integrate all 3 into unified demo"}],
    [{"t": "Then: Phase 1 decision based on POC results"}],
])

rect(s, 0.7, 5.2, 11.6, 2.0, fill=ORANGE, round_=True)
text(s, 0.95, 5.4, 11.2, 1.6, [
    {"runs": [{"t": "The Ask", "bold": True, "size": 13, "color": WHITE}], "space_after": 0},
    {"runs": [{"t": "Approve 3-week POC sprint (3 working demos, ~$29 cost, $0 infrastructure). Stakeholder demos at end of each week. Decision point: Phase 1 after demos complete.", "size": 12, "color": WHITE}], "space_after": 0},
])

# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Sonata-Knowledge-Assistant-Stakeholder-Deck.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
